"""
Blueprint para gestión de informes del fiscalizador - registrar, listar,
editar, eliminar, descargar y analizar informes con IA
"""

from flask import Blueprint, render_template, request, redirect, session, send_file, url_for, jsonify, flash
from werkzeug.utils import secure_filename
from conexion import conexion
from utils import error_response, acceso_no_autorizado, error_interno, datos_invalidos, no_encontrado
from utils.validators import archivo_permitido, sanitizar_nombre, validar_mime_real, validar_longitudes
from datetime import datetime
from io import BytesIO
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import landscape, letter
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import json
import re
from dotenv import load_dotenv

# Cargar variables de entorno desde archivo .env
load_dotenv()

# Blueprint principal de informes
informe_bp = Blueprint("informe", __name__)

# Configuración de subida de archivos
UPLOAD_FOLDER = "uploads"
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB

# Crear carpeta de subidas si no existe
os.makedirs(UPLOAD_FOLDER, exist_ok=True)


# ==========================================
# LISTAR INFORMES DEL FISCALIZADOR
# ==========================================
@informe_bp.route("/mis_informes")
def mis_informes():

    if "usuario" not in session:
        return redirect(url_for("home"))

    if session.get("perfil_activo") != "fiscalizador":
        return acceso_no_autorizado()

    fecha_desde = request.args.get("fecha_desde")
    fecha_hasta = request.args.get("fecha_hasta")
    tipo = request.args.get("tipo")
    anio = request.args.get("anio")
    titulo = request.args.get("titulo")

    conn = None
    cursor = None

    try:
        conn = conexion()
        cursor = conn.cursor(dictionary=True)

        sql = """
            SELECT
                i.*,
                g.fecha_guardia
            FROM informes i
            INNER JOIN guardias g
                ON i.id_guardia = g.id_guardia
            WHERE i.id_usuario = %s
            AND i.estado = 'activo'
        """

        parametros = [session["id_usuario"]]

        # ── Construcción dinámica de la consulta con filtros opcionales ──
        if fecha_desde:
            sql += " AND g.fecha_guardia >= %s "
            parametros.append(fecha_desde)

        # FILTRO FECHA HASTA
        if fecha_hasta:
            sql += " AND g.fecha_guardia <= %s "
            parametros.append(fecha_hasta)

        # FILTRO TIPO
        if tipo:
            sql += " AND i.tipo_archivo = %s "
            parametros.append(tipo)

        # FILTRO AÑO
        if anio:
            sql += " AND YEAR(g.fecha_guardia) = %s "
            parametros.append(anio)

        # FILTRO TITULO
        if titulo:
            sql += " AND i.titulo LIKE %s "
            parametros.append(f"%{titulo}%")

        sql += " ORDER BY i.fecha_subida DESC "

        cursor.execute(sql, tuple(parametros))
        informes = cursor.fetchall()

        anio_actual = datetime.now().year

        return render_template(
            "mis_informes.html",
            informes=informes,
            anio_actual=anio_actual
        )

    finally:
        if cursor is not None: cursor.close()
        if conn is not None: conn.close()


# ==========================================
# REGISTRAR INFORME
# ==========================================
@informe_bp.route("/registrar_informe", methods=["GET", "POST"])
def registrar_informe():

    if "usuario" not in session:
        return redirect(url_for("home"))

    if session.get("perfil_activo") != "fiscalizador":
        return acceso_no_autorizado()

    conn = None
    cursor = None

    try:
        conn = conexion()
        cursor = conn.cursor(dictionary=True)

        if request.method == "POST":

            id_guardia = request.form["id_guardia"]
            titulo = request.form["titulo"]
            descripcion = request.form["descripcion"]

            valido, msg = validar_longitudes({
                "titulo": titulo,
                "descripcion": descripcion,
            })
            if not valido:
                return datos_invalidos(msg)

            archivo = request.files.get("archivo")

            if not archivo or archivo.filename == "":
                return datos_invalidos("Debe seleccionar un archivo")

            if not archivo_permitido(archivo.filename):
                return datos_invalidos("Formato de archivo no permitido")

            extension = archivo.filename.rsplit(".", 1)[1].lower()

            # ── Validar tamaño máximo ──
            archivo.seek(0, os.SEEK_END)
            tamano_archivo = archivo.tell()
            archivo.seek(0)
            if tamano_archivo > MAX_FILE_SIZE:
                return datos_invalidos("El archivo excede el tamaño máximo permitido (10 MB)")

            # ── Validar MIME real ──
            magic_bytes = archivo.read(12)
            archivo.seek(0)
            if not validar_mime_real(magic_bytes, extension):
                return datos_invalidos("El contenido del archivo no coincide con su extensión")

            nombre = (
                datetime.now().strftime("%Y%m%d%H%M%S_")
                + sanitizar_nombre(secure_filename(archivo.filename))
            )

            ruta = os.path.join(
                UPLOAD_FOLDER,
                nombre
            )

            archivo.save(ruta)

            cursor.execute("""
                INSERT INTO informes(
                    id_guardia,
                    id_usuario,
                    titulo,
                    descripcion,
                    nombre_archivo,
                    ruta_archivo,
                    tipo_archivo,
                    extension,
                    tamano_archivo
                )
                VALUES(
                    %s,%s,%s,%s,
                    %s,%s,%s,%s,%s
                )
            """, (
                id_guardia,
                session["id_usuario"],
                titulo,
                descripcion,
                nombre,
                ruta,
                extension,
                extension,
                tamano_archivo
            ))

            conn.commit()

            flash("Informe registrado correctamente", "success")
            return redirect(url_for("informe.mis_informes"))

        cursor.execute("""
            SELECT
                id_guardia,
                fecha_guardia
            FROM guardias
            WHERE id_usuario = %s
            ORDER BY fecha_guardia DESC
        """, (session["id_usuario"],))

        guardias = cursor.fetchall()

        return render_template(
            "registrar_informe.html",
            guardias=guardias
        )

    finally:
        if cursor is not None: cursor.close()
        if conn is not None: conn.close()


# ==========================================
# EDITAR INFORME (título, descripción y archivo)
# ==========================================
@informe_bp.route("/editar_informe/<int:id_informe>", methods=["GET", "POST"])
def editar_informe(id_informe):

    if "usuario" not in session:
        return redirect(url_for("home"))

    if session.get("perfil_activo") != "fiscalizador":
        return acceso_no_autorizado()

    conn = None
    cursor = None

    try:
        conn = conexion()
        cursor = conn.cursor(dictionary=True)

        cursor.execute("""
            SELECT *
            FROM informes
            WHERE id_informe = %s
            AND id_usuario = %s
            AND estado = 'activo'
        """, (
            id_informe,
            session["id_usuario"]
        ))

        informe = cursor.fetchone()

        if not informe:
            return no_encontrado("Informe no encontrado")

        if request.method == "POST":

            titulo = request.form["titulo"]
            descripcion = request.form["descripcion"]

            valido, msg = validar_longitudes({
                "titulo": titulo,
                "descripcion": descripcion,
            })
            if not valido:
                return datos_invalidos(msg)

            cursor.execute("""
                UPDATE informes
                SET titulo=%s,
                    descripcion=%s
                WHERE id_informe=%s
            """, (
                titulo,
                descripcion,
                id_informe
            ))

            archivo = request.files.get("archivo")

            if archivo and archivo.filename != "":

                if archivo_permitido(archivo.filename):

                    extension_edit = archivo.filename.rsplit(".", 1)[1].lower()

                    # ── Validar tamaño máximo ──
                    archivo.seek(0, os.SEEK_END)
                    tamano_nuevo = archivo.tell()
                    archivo.seek(0)
                    if tamano_nuevo > MAX_FILE_SIZE:
                        return datos_invalidos("El archivo excede el tamaño máximo permitido (10 MB)")

                    # ── Validar MIME real ──
                    magic_bytes = archivo.read(12)
                    archivo.seek(0)
                    if not validar_mime_real(magic_bytes, extension_edit):
                        return datos_invalidos("El contenido del archivo no coincide con su extensión")

                    nombre = (
                        datetime.now().strftime("%Y%m%d%H%M%S_")
                        + sanitizar_nombre(secure_filename(archivo.filename))
                    )

                    ruta = os.path.join(
                        UPLOAD_FOLDER,
                        nombre
                    )

                    archivo.save(ruta)

                    extension = nombre.rsplit(".", 1)[1].lower()

                    cursor.execute("""
                        UPDATE informes
                        SET nombre_archivo=%s,
                            ruta_archivo=%s,
                            tipo_archivo=%s,
                            extension=%s,
                            tamano_archivo=%s
                        WHERE id_informe=%s
                    """, (
                        nombre,
                        ruta,
                        extension,
                        extension,
                        tamano_nuevo,
                        id_informe
                    ))

            conn.commit()

            flash("Informe actualizado correctamente", "success")
            return redirect(url_for("informe.mis_informes"))

        return render_template(
            "editar_informe.html",
            informe=informe
        )

    finally:
        if cursor is not None: cursor.close()
        if conn is not None: conn.close()


# ==========================================
# DESCARGAR ARCHIVO DEL INFORME
# ==========================================
@informe_bp.route("/descargar_informe/<int:id_informe>")
def descargar_informe(id_informe):

    if "usuario" not in session:
        return redirect(url_for("home"))

    conn = None
    cursor = None

    try:
        conn = conexion()
        cursor = conn.cursor(dictionary=True)

        cursor.execute("""
            SELECT *
            FROM informes
            WHERE id_informe = %s
            AND estado = 'activo'
        """, (id_informe,))

        informe = cursor.fetchone()

        if not informe:
            flash("Archivo no encontrado", "error")
            return redirect(url_for("informe.mis_informes"))

        ruta = informe["ruta_archivo"]
        # ── Protección contra Path Traversal ──
        ruta_real = os.path.realpath(ruta)
        uploads_real = os.path.realpath(UPLOAD_FOLDER)
        if not ruta_real.startswith(uploads_real + os.sep) and ruta_real != uploads_real:
            return acceso_no_autorizado()

        if not os.path.exists(ruta_real):
            flash("El archivo no existe en el servidor", "error")
            return redirect(url_for("informe.mis_informes"))

        return send_file(
            ruta_real,
            as_attachment=True,
            download_name=informe["nombre_archivo"]
        )

    finally:
        if cursor is not None: cursor.close()
        if conn is not None: conn.close()


# ==========================================
# ELIMINAR INFORME (SOFT DELETE)
# ==========================================
@informe_bp.route("/eliminar_informe/<int:id_informe>", methods=["POST"])
def eliminar_informe(id_informe):

    if "usuario" not in session:
        return redirect(url_for("home"))

    if session.get("perfil_activo") != "fiscalizador":
        return acceso_no_autorizado()

    conn = None
    cursor = None

    try:
        conn = conexion()
        cursor = conn.cursor()

        cursor.execute("""
            UPDATE informes
            SET estado = 'eliminado'
            WHERE id_informe = %s
            AND id_usuario = %s
        """, (id_informe, session["id_usuario"]))

        conn.commit()

        flash("Informe eliminado correctamente", "success")
        return redirect(url_for("informe.mis_informes"))

    finally:
        if cursor is not None: cursor.close()
        if conn is not None: conn.close()


# ==========================================
# EXPORTAR LISTADO A PDF
# ==========================================
@informe_bp.route("/mis_informes/exportar/pdf")
def exportar_informes_pdf():

    if "usuario" not in session:
        return redirect(url_for("home"))

    if session.get("perfil_activo") != "fiscalizador":
        return acceso_no_autorizado()

    fecha_desde = request.args.get("fecha_desde")
    fecha_hasta = request.args.get("fecha_hasta")
    tipo = request.args.get("tipo")
    anio = request.args.get("anio")
    titulo = request.args.get("titulo")

    conn = None
    cursor = None

    try:
        conn = conexion()
        cursor = conn.cursor(dictionary=True)

        sql = """
            SELECT
                i.titulo,
                i.descripcion,
                i.nombre_archivo,
                i.tipo_archivo,
                i.fecha_subida,
                g.fecha_guardia
            FROM informes i
            INNER JOIN guardias g
                ON i.id_guardia = g.id_guardia
            WHERE i.id_usuario = %s
            AND i.estado = 'activo'
        """

        parametros = [session["id_usuario"]]

        if fecha_desde:
            sql += " AND g.fecha_guardia >= %s "
            parametros.append(fecha_desde)

        if fecha_hasta:
            sql += " AND g.fecha_guardia <= %s "
            parametros.append(fecha_hasta)

        if tipo:
            sql += " AND i.tipo_archivo = %s "
            parametros.append(tipo)

        if anio:
            sql += " AND YEAR(g.fecha_guardia) = %s "
            parametros.append(anio)

        if titulo:
            sql += " AND i.titulo LIKE %s "
            parametros.append(f"%{titulo}%")

        sql += " ORDER BY i.fecha_subida DESC "

        cursor.execute(sql, tuple(parametros))
        data = cursor.fetchall()

        buffer = BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=landscape(letter))
        styles = getSampleStyleSheet()
        elementos = []

        elementos.append(Paragraph("MIS INFORMES", styles["Title"]))
        elementos.append(Spacer(1, 10))

        table_data = [["Fecha Guardia", "Título", "Descripción", "Archivo", "Tipo", "Registro"]]

        for d in data:
            table_data.append([
                str(d["fecha_guardia"]) if d["fecha_guardia"] else "",
                d["titulo"] or "",
                (d["descripcion"] or "")[:80],
                d["nombre_archivo"] or "",
                (d["tipo_archivo"] or "").upper(),
                str(d["fecha_subida"]) if d["fecha_subida"] else ""
            ])

        table = Table(table_data, repeatRows=1)

        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.darkblue),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
            ("ALIGN", (0, 0), (-1, -1), "CENTER"),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("BACKGROUND", (0, 1), (-1, -1), colors.white),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.lightgrey]),
        ]))

        elementos.append(table)

        doc.build(elementos)
        buffer.seek(0)

        return send_file(buffer, download_name="mis_informes.pdf", as_attachment=True)

    finally:
        if cursor is not None: cursor.close()
        if conn is not None: conn.close()


# ==========================================
# EXPORTAR LISTADO A EXCEL
# ==========================================
@informe_bp.route("/mis_informes/exportar/excel")
def exportar_informes_excel():

    if "usuario" not in session:
        return redirect(url_for("home"))

    if session.get("perfil_activo") != "fiscalizador":
        return acceso_no_autorizado()

    fecha_desde = request.args.get("fecha_desde")
    fecha_hasta = request.args.get("fecha_hasta")
    tipo = request.args.get("tipo")
    anio = request.args.get("anio")
    titulo = request.args.get("titulo")

    conn = None
    cursor = None

    try:
        conn = conexion()
        cursor = conn.cursor(dictionary=True)

        sql = """
            SELECT
                i.titulo,
                i.descripcion,
                i.nombre_archivo,
                i.tipo_archivo,
                i.fecha_subida,
                g.fecha_guardia
            FROM informes i
            INNER JOIN guardias g
                ON i.id_guardia = g.id_guardia
            WHERE i.id_usuario = %s
            AND i.estado = 'activo'
        """

        parametros = [session["id_usuario"]]

        if fecha_desde:
            sql += " AND g.fecha_guardia >= %s "
            parametros.append(fecha_desde)

        if fecha_hasta:
            sql += " AND g.fecha_guardia <= %s "
            parametros.append(fecha_hasta)

        if tipo:
            sql += " AND i.tipo_archivo = %s "
            parametros.append(tipo)

        if anio:
            sql += " AND YEAR(g.fecha_guardia) = %s "
            parametros.append(anio)

        if titulo:
            sql += " AND i.titulo LIKE %s "
            parametros.append(f"%{titulo}%")

        sql += " ORDER BY i.fecha_subida DESC "

        cursor.execute(sql, tuple(parametros))
        data = cursor.fetchall()

        wb = Workbook()
        ws = wb.active
        ws.title = "Mis Informes"

        header_font = Font(name="Segoe UI", bold=True, color="FFFFFF", size=11)
        header_fill = PatternFill(start_color="1E40AF", end_color="1E40AF", fill_type="solid")
        header_alignment = Alignment(horizontal="center", vertical="center")
        thin_border = Border(
            left=Side(style="thin"),
            right=Side(style="thin"),
            top=Side(style="thin"),
            bottom=Side(style="thin")
        )

        headers = ["Fecha Guardia", "Título", "Descripción", "Archivo", "Tipo", "Registro"]
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = header_alignment
            cell.border = thin_border

        for row_idx, d in enumerate(data, 2):
            row_data = [
                str(d["fecha_guardia"]) if d["fecha_guardia"] else "",
                d["titulo"] or "",
                d["descripcion"] or "",
                d["nombre_archivo"] or "",
                (d["tipo_archivo"] or "").upper(),
                str(d["fecha_subida"]) if d["fecha_subida"] else ""
            ]
            for col_idx, val in enumerate(row_data, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=val)
                cell.border = thin_border
                cell.alignment = Alignment(vertical="center")

        ws.column_dimensions["A"].width = 15
        ws.column_dimensions["B"].width = 30
        ws.column_dimensions["C"].width = 40
        ws.column_dimensions["D"].width = 35
        ws.column_dimensions["E"].width = 12
        ws.column_dimensions["F"].width = 20

        output = BytesIO()
        wb.save(output)
        output.seek(0)

        return send_file(output, download_name="mis_informes.xlsx", as_attachment=True)

    finally:
        if cursor is not None: cursor.close()
        if conn is not None: conn.close()


# ==========================================
# FUNCIONES AUXILIARES PARA ANÁLISIS CON IA
# ==========================================

def _extraer_texto_archivo(archivo_obj):
    """Extrae texto de PDF, DOCX o XLSX y lo devuelve como string."""
    ext = archivo_obj["extension"].lower()
    ruta = archivo_obj["ruta_archivo"]

    texto = ""

    if ext == "pdf":
        try:
            import pdfplumber

            def _es_amarillo(color):
                try:
                    if isinstance(color, (tuple, list)):
                        if len(color) == 3:
                            r, g, b = color
                            return r >= 0.75 and g >= 0.75 and b <= 0.45
                        if len(color) == 4:
                            c, m, y, k = color
                            return c <= 0.25 and m <= 0.35 and y >= 0.6 and k <= 0.25
                except Exception:
                    pass
                return False

            with pdfplumber.open(ruta) as pdf:
                paginas = []
                for num, pagina in enumerate(pdf.pages, 1):
                    rects_amarillos = []
                    try:
                        for r in pagina.rects:
                            if _es_amarillo(r.get("non_stroking_color")):
                                rects_amarillos.append((r["x0"], r["top"], r["x1"], r["bottom"]))
                    except Exception:
                        pass

                    lineas = []
                    try:
                        for lin in pagina.extract_text_lines():
                            t = (lin.get("text") or "").strip()
                            if not t:
                                continue
                            marca = ""
                            if rects_amarillos:
                                for (x0, top, x1, bottom) in rects_amarillos:
                                    if lin["top"] < bottom and lin["bottom"] > top and lin["x0"] < x1 and lin["x1"] > x0:
                                        marca = "[CRITICO-AMARILLO] "
                                        break
                            lineas.append(marca + t)
                    except Exception:
                        t = pagina.extract_text()
                        if t:
                            lineas.append(t)
                    if lineas:
                        paginas.append(f"===== PÁGINA: {num} =====\n" + "\n".join(lineas))
                texto = "\n".join(paginas)
        except Exception as e:
            texto = f"[Error al extraer texto del PDF: {str(e)}]"

    elif ext == "docx":
        try:
            from docx import Document
            from docx.enum.text import WD_COLOR_INDEX

            def _parrafo_amarillo(p):
                try:
                    return any(r.font.highlight_color == WD_COLOR_INDEX.YELLOW for r in p.runs)
                except Exception:
                    return False

            def _celda_amarilla(celda):
                try:
                    tcPr = celda._tc.tcPr
                    if tcPr is not None:
                        shd = tcPr.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}shd")
                        if shd is not None:
                            fill = (shd.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fill") or "").upper()
                            if fill == "FFFF00" or (len(fill) == 6 and fill not in ("AUTO", "FFFFFF") and fill.startswith("FF") and fill.endswith("00")):
                                return True
                    return any(_parrafo_amarillo(p) for p in celda.paragraphs)
                except Exception:
                    return False

            doc = Document(ruta)
            parrafos = []
            for p in doc.paragraphs:
                if p.text.strip():
                    marca = "[CRITICO-AMARILLO] " if _parrafo_amarillo(p) else ""
                    parrafos.append(marca + p.text)
            texto = "\n".join(parrafos)

            # Intentar extraer tablas del DOCX
            if doc.tables:
                texto += "\n\n--- TABLAS ENCONTRADAS ---\n"
                for idx, tabla in enumerate(doc.tables, 1):
                    texto += f"\nTabla {idx}:\n"
                    for fila in tabla.rows:
                        celdas = []
                        fila_amarilla = False
                        for celda in fila.cells:
                            celdas.append(celda.text.strip().replace("\n", " / "))
                            if _celda_amarilla(celda):
                                fila_amarilla = True
                        if any(c for c in celdas):
                            marca = "[CRITICO-AMARILLO] " if fila_amarilla else ""
                            texto += marca + " | ".join(celdas) + "\n"
        except Exception as e:
            texto = f"[Error al extraer texto del DOCX: {str(e)}]"

    elif ext in ("xlsx", "xlsm"):
        try:
            from openpyxl import load_workbook
            wb = load_workbook(ruta, data_only=True)
            for nombre_hoja in wb.sheetnames:
                ws = wb[nombre_hoja]
                texto += f"\n===== HOJA: {nombre_hoja} =====\n"
                for fila in ws.iter_rows():
                    celdas = []
                    fila_amarilla = False
                    for c in fila:
                        valor = str(c.value).replace("\n", " / ") if c.value is not None else ""
                        celdas.append(valor)
                        try:
                            f = c.fill
                            if f is not None and f.patternType == "solid":
                                rgb = getattr(f.fgColor, "rgb", None)
                                if isinstance(rgb, str) and rgb.upper() in ("FFFFFF00", "00FFFF00"):
                                    fila_amarilla = True
                        except Exception:
                            pass
                    if any(c for c in celdas):
                        marca = "[CRITICO-AMARILLO] " if fila_amarilla else ""
                        texto += marca + " | ".join(celdas) + "\n"
        except Exception as e:
            texto = f"[Error al extraer texto del XLSX: {str(e)}]"

    return texto.strip()


def _validar_coherencia_analisis(resultado):
    """Corrige de forma determinista las inconsistencias aritméticas del análisis:
    los totales del módulo 1 deben coincidir con los desgloses y los módulos 2, 3 y 4."""

    # Ajustar conteos: trabajos críticos, producción recuperada y total de actividades
    try:
        m1 = resultado.get("modulo1_resumen_ejecutivo")
        if not isinstance(m1, dict):
            return resultado

        # Trabajos críticos = cantidad real de filas en modulo3
        criticos = resultado.get("modulo3_trabajos_importantes")
        if isinstance(criticos, list):
            m1["total_trabajos_criticos"] = len(criticos)

        # Producción recuperada = suma exacta de bopd_por_taller
        m4 = resultado.get("modulo4_resumen_operativo")
        if isinstance(m4, dict):
            bopd = m4.get("bopd_por_taller")
            if isinstance(bopd, dict) and bopd:
                valores = [v for v in bopd.values() if isinstance(v, (int, float))]
                m1["total_produccion_recuperada_bopd"] = round(sum(valores), 2)
            tipos = m4.get("trabajos_por_tipo")
            if isinstance(tipos, dict) and tipos:
                m1["desglose_por_tipo"] = tipos

        # Total de actividades = suma de ejecutadas del desglose por taller
        desglose = m1.get("desglose_por_taller")
        if isinstance(desglose, dict) and desglose:
            total_ej = 0
            valido = True
            for v in desglose.values():
                if isinstance(v, dict) and isinstance(v.get("ejecutadas"), (int, float)):
                    total_ej += v["ejecutadas"]
                elif isinstance(v, (int, float)):
                    total_ej += v
                else:
                    valido = False
                    break
            if valido and total_ej > 0:
                m1["total_actividades"] = int(total_ej)

        # Eficacia por taller recalculada exactamente
        m2 = resultado.get("modulo2_eficacia_taller")
        if isinstance(m2, list):
            for fila in m2:
                if isinstance(fila, dict):
                    ej = fila.get("ejecutadas")
                    comp = fila.get("completadas")
                    if isinstance(ej, (int, float)) and isinstance(comp, (int, float)) and ej > 0:
                        fila["eficacia"] = round(comp * 100.0 / ej, 2)
    except Exception:
        pass
    return resultado


def _analizar_con_gemini(texto, titulo, descripcion):
    """Envía el texto a Gemini y devuelve la respuesta estructurada.
    Prueba múltiples modelos gratuitos en cascada si hay error de cuota."""

    # Configuración de API Key y modelos de fallback
    from google import genai
    import time

    api_key = os.getenv("GEMINI_API_KEY")
    model_principal = os.getenv("GEMINI_MODEL", "gemini-2.5-flash")

    if not api_key or api_key == "TU_GEMINI_API_KEY":
        return {
            "error": "API Key de Gemini no configurada. Agrega GEMINI_API_KEY en el archivo .env",
            "recomendaciones": []
        }

    # Modelos gratuitos de Google en orden de prioridad (más capaz → menos capaz)
    # Se prueban en cascada: si uno se queda sin cuota, el siguiente toma el relevo
    modelos_fallback = [
        model_principal,
        "gemini-2.5-flash-lite",
        "gemini-2.0-flash",
        "gemini-2.0-flash-lite",
    ]
    # Eliminar duplicados manteniendo el orden
    modelos = list(dict.fromkeys(modelos_fallback))

    client = genai.Client(api_key=api_key)

    prompt = f"""ROL Y CONTEXTO
Actúas como un Motor de Procesamiento de Datos Operativos y Analista Senior de Producción en Petróleo y Gas. Tu tarea es analizar con absoluta precisión matemática el siguiente archivo de reporte de guardia y estructurar un análisis analítico detallado. El resultado final debe ser un informe operativo riguroso, consistente y libre de errores matemáticos que sirva como fuente única de verdad sobre los eventos de la guardia.
Debes ceñirte estrictamente a las reglas de negocio, lógica de cálculo e instrucciones de mapeo que se detallan a continuación.

TÍTULO DEL DOCUMENTO: {titulo}
DESCRIPCIÓN: {descripcion}

CONTENIDO DEL DOCUMENTO:
{texto[:30000]}

1. ARQUITECTURA LÓGICA DEL DOCUMENTO (INDEPENDIENTE DEL FORMATO)
El documento original puede ser Excel, PDF o Word, pero SIEMPRE sigue el mismo modelo de estructura (los nombres, fechas y cantidades son variables; adáptate SIEMPRE a lo que contenga ESTE documento):
- El documento se segmenta por días de registro diario (pueden ser 1, 2, 3 o más días). Si proviene de Excel, cada día viene como una hoja marcada con "===== HOJA: <nombre> =====" (ej. "===== HOJA: 27 de Junio ====="); si proviene de PDF verás marcas "===== PÁGINA: n ====="; si proviene de Word el texto es continuo. En todos los casos, el inicio de cada día se identifica por un título tipo "TRABAJOS DE GUARDIA: <DÍA> <FECHA>". Usa las fechas y días REALES del documento.
- Puede existir una sección u hoja final dedicada a tareas en espera llamada "Pendientes" o "TRABAJOS PENDIENTES". Si no existe, los pendientes se determinan solo por los estados de las actividades diarias.
- La información no es una sola tabla continua. Está segmentada verticalmente por bloques de Taller.
- Identificador de Taller: El inicio de un bloque está marcado por una fila o línea cuyo único contenido es el nombre del taller en mayúsculas. Ejemplos típicos (NO exhaustivos, usa los que aparezcan en el documento): "MONTAJE", "MONTAJE - ADICIONAL", "MECÁNICOS", "ENERGÍA", "INSTRUMENTACIÓN", "COMPRESIÓN DE GAS", "GASFITERÍA Y SOLDADURA", "MOVIMIENTO DE SUELOS". El documento puede contener otros talleres distintos: inclúyelos todos tal como aparezcan.
- Encabezados de Bloque típicos: Ítem | Pozo | Batería | Producción (bopd) | Requerimiento | Estado | Fecha de Ejecución | Tipo | Actividad Ejecutada. Pueden variar ligeramente u omitirse columnas (ej. la sección "Pendientes" puede no tener columna Tipo): usa SIEMPRE la fila de encabezados real de cada bloque.
- Un bloque de taller finaliza cuando se inicia un taller diferente o una nueva hoja/sección.
- Las filas de tablas vienen con columnas separadas por " | " (Excel y tablas de Word). En PDF las columnas pueden venir separadas solo por espacios: alinéalas igualmente con su encabezado. La primera posición puede venir vacía: NO cuentes posiciones fijas, alinea cada valor con su encabezado de columna correspondiente.
- Las filas o líneas marcadas con el prefijo "[CRITICO-AMARILLO]" corresponden a celdas o texto resaltado en amarillo en el documento original: son los Trabajos Importantes/Críticos.

2. MAPEO EXACTO DE CAMPOS
Para cada registro, usa la fila de encabezados de su bloque para mapear cada valor con su columna (NO uses posiciones fijas):
- Taller Evaluado: Título del taller activo que encabeza el bloque.
- Pozo: Código identificador del pozo/equipo (columna "Pozo"). Puede ser un pozo, o un equipo/ubicación (ej. BAT, ESTACION, MANIFOLD, PLANTA).
- Producción (bopd): Volumen numérico (columna "Producción (bopd)"). Puede venir con coma decimal (ej. "3,10" = 3.10) o "-"/vacío (= sin dato, trátalo como 0 para sumas).
- Requerimiento (Falla): Descripción del problema técnico (columna "Requerimiento").
- Estado: Nivel de completación (columna "Estado"). El valor 1 equivale a 100% completado. Un valor decimal entre 0 y 1 (ej. 0.2, 0.75) indica avance parcial = "en proceso". 0, texto o vacío significa "pendiente".
- Tipo de Actividad: Categoría del trabajo (columna "Tipo"). Son códigos cortos (ej. "CNP", "PV", "SOP", "SUS", "MC", "OP"). Usa SOLO los códigos que aparezcan en ESTE documento, sean cuales sean.
- Actividad Ejecutada (Solución): Descripción de la solución aplicada en campo (columna "Actividad Ejecutada"). Las líneas internas de la celda vienen separadas por " / ".

3. ALGORITMO DE LIMPIEZA Y CONSOLIDACIÓN (REGLAS DE NEGOCIO OBLIGATORIAS)
- Regla 1: Unificación de Talleres con la Misma Raíz: Consolida sub-áreas bajo una sola categoría principal. Aplica esto a CUALQUIER taller cuyo nombre tenga la misma raíz más un sufijo como "ADICIONAL", "- ADICIONAL" o similar. Ej: "MONTAJE" y "MONTAJE - ADICIONAL" se unifican como "MONTAJE"; "ENERGÍA" y "ENERGÍA ADICIONAL" como "ENERGÍA".
- Regla 2: Deduplicación Temporal: Si una misma actividad (mismo pozo y misma descripción de solución o requerimiento) se registra en varios días (consecutivos o no), contabilízala como 1 sola actividad en el consolidado, tomando el Estado del último día reportado.
- Regla 3: Deduplicación Inter-Taller: Si la misma actividad exacta en el mismo pozo es reportada bajo dos talleres diferentes (o en el bloque normal y su bloque ADICIONAL del mismo día), contabilízala como 1 sola actividad, asignándosela al taller que reportó la ejecución principal o cierre.
- Regla 4: Filas "Sin novedad" o completamente vacías (sin Pozo ni Requerimiento) NO cuentan como actividades.

4. INSTRUCCIONES DE CÁLCULO DE MÉTRICAS Y KPIs
- Conteo de Actividades Totales vs. Completadas: Suma actividades únicas por taller unificado (aplicando Reglas 1, 2, 3 y 4). Cuenta cuántas tienen Estado = 1. Eficacia (%) = (Completadas / Totales) * 100.
- Volumen de Producción Recuperada (CNP): Filtra registros ÚNICOS (ya deduplicados) donde Tipo = "CNP" Y Estado = 1. Suma valores de Producción (bopd) (convierte comas decimales a punto; "-" o vacío = 0). Reporta total general y desglose por taller unificado.
- Identificación de Trabajos Importantes (Críticos): Son EXCLUSIVAMENTE las filas o líneas marcadas con el prefijo "[CRITICO-AMARILLO]" (resaltado amarillo en el documento original). PROHIBIDO incluir cualquier actividad que no tenga ese prefijo, aunque parezca importante o urgente. Cuenta cantidad por taller y construye matriz: Taller | Pozo | Falla Detectada | Solución Operativa. Si no hay ninguna fila marcada, reporta 0 trabajos críticos y deja la matriz vacía.
- Clasificación de Trabajos por Tipo: Agrupa y cuenta tareas únicas por código de tipo, usando SOLO los códigos que realmente aparecen en el documento (ej. CNP, PV, SOP, SUS).
- Cierre de Guardia (Pendientes): Procedimiento OBLIGATORIO para determinar el estado FINAL de cada tarea:
  a) Agrupa las apariciones de la misma tarea (mismo pozo + mismo requerimiento o equivalente) a través de TODOS los días.
  b) El estado final es el del ÚLTIMO día donde aparece. Ejemplo: si una tarea aparece un día con Estado 0 y un día posterior con Estado 1, la tarea está COMPLETADA y NO va en modulo5_pendientes.
  c) modulo5_pendientes incluye SOLO tareas con estado final distinto de 1: usa "en proceso" si el estado final es decimal entre 0 y 1 (ej. 0.2, 0.8), y "pendiente" si es 0, vacío o texto.
  d) Si existe una hoja o sección "Pendientes", incluye también sus filas con datos reales (ignora filas sin Pozo ni Requerimiento). Clasifica todo por taller unificado.

5. RESTRICCIÓN ESTRICTA DE FIDELIDAD DE DATOS (CERO ALUCINACIONES)
- No inventes, asumas ni extrapoles ningún pozo, código, volumen, falla, solución o tarea que no esté explícitamente en el documento.
- Si una celda o campo está vacío, repórtalo como "No especificado" o "Vacío".
- Todos los totales, promedios y porcentajes deben ser el resultado exacto de la sumatoria de los datos del documento.

6. ESTRUCTURA DE SALIDA REQUERIDA
Devuelve ÚNICAMENTE un JSON válido (sin markdown, sin comillas triples) con esta estructura exacta:

{{
  "resumen": "Pega aquí el texto completo del MÓDULO 1 (Resumen Ejecutivo) en formato legible, incluyendo total de actividades, producción recuperada en BOPD, total de trabajos críticos y desglose por tipo y taller.",
  "hallazgos": ["Hallazgo clave 1", "Hallazgo clave 2", "Hallazgo clave 3", "Hallazgo clave 4", "Hallazgo clave 5"],
  "recomendaciones": ["Recomendación 1", "Recomendación 2", "Recomendación 3", "Recomendación 4", "Recomendación 5"],
  "graficas": [
    {{
      "tipo": "bar",
      "titulo": "Actividades por Taller Unificado",
      "labels": ["<TALLERES REALES UNIFICADOS DEL DOCUMENTO>"],
      "datasets": [
        {{"label": "Ejecutadas", "data": [0]}},
        {{"label": "Completadas", "data": [0]}}
      ]
    }},
    {{
      "tipo": "pie",
      "titulo": "Distribución de Trabajos por Tipo",
      "labels": ["<TIPOS REALES DEL DOCUMENTO, ej. CNP, PV, SOP, SUS>"],
      "datasets": [
        {{"label": "Cantidad", "data": [0], "backgroundColor": ["#22c55e", "#f59e0b", "#3b82f6", "#8b5cf6"]}}
      ]
    }},
    {{
      "tipo": "bar",
      "titulo": "Producción Recuperada CNP por Taller (BOPD)",
      "labels": ["<TALLERES REALES CON CNP COMPLETADO>"],
      "datasets": [
        {{"label": "BOPD Recuperados", "data": [0], "backgroundColor": ["#1e3a5f", "#2563eb", "#3b82f6", "#60a5fa"]}}
      ]
    }}
  ],
  "modulo1_resumen_ejecutivo": {{
    "total_actividades": 0,
    "total_produccion_recuperada_bopd": 0,
    "total_trabajos_criticos": 0,
    "desglose_por_tipo": {{}},
    "desglose_por_taller": {{}}
  }},
  "modulo2_eficacia_taller": [
    {{"taller": "MONTAJE", "ejecutadas": 0, "completadas": 0, "eficacia": 0.0}}
  ],
  "modulo3_trabajos_importantes": [
    {{"taller": "MONTAJE", "pozo": "XXX-001", "falla": "Descripción de la falla", "solucion": "Descripción de la solución"}}
  ],
  "modulo4_resumen_operativo": {{
    "bopd_por_taller": {{}},
    "trabajos_por_tipo": {{}}
  }},
  "modulo5_pendientes": [
    {{"taller": "MONTAJE", "pozo": "XXX-001", "requerimiento": "Descripción de la tarea pendiente", "estado": "pendiente"}}
  ]
}}

IMPORTANTE: 
- Responde SOLO con el JSON válido, sin explicaciones ni markdown.
- Reemplaza TODOS los valores de ejemplo con los datos REALES extraídos del documento.
- Llena TODOS los arrays y objetos con los datos correctos según tu análisis.
- El campo "resumen" debe contener el texto completo del MÓDULO 1: RESUMEN EJECUTIVO DE GUARDIA con formato legible (usa saltos de línea \\n para separar secciones).
- VERIFICACIÓN FINAL OBLIGATORIA antes de responder:
  a) modulo3_trabajos_importantes debe contener TODAS las filas marcadas "[CRITICO-AMARILLO]" (tras deduplicar por pozo+requerimiento) y NINGUNA actividad sin esa marca.
  b) modulo5_pendientes SOLO debe contener tareas cuyo Estado FINAL (último día donde aparece la tarea) sea distinto de 1. Revisa tarea por tarea: si su última aparición tiene Estado = 1, elimínala de pendientes (ej. una tarea con Estado 1 en todas sus apariciones NUNCA es pendiente). Y toda tarea cuya última aparición tenga Estado 0, decimal o vacío DEBE estar en pendientes, aunque aparezca un solo día.
  c) Recalcula cada suma: total_produccion_recuperada_bopd debe ser EXACTAMENTE igual a la suma de bopd_por_taller, y total_actividades igual a la suma de "ejecutadas" del desglose_por_taller. Corrige cualquier inconsistencia antes de responder.
  d) bopd_por_taller debe incluir SOLO talleres cuya suma de CNP completados sea mayor que 0 (omite talleres con 0)."""

    max_intentos_por_modelo = 2
    ultimo_error = None
    texto_respuesta = ""

    for modelo in modelos:
        for intento in range(1, max_intentos_por_modelo + 1):
            try:
                respuesta = client.models.generate_content(
                    model=modelo,
                    contents=prompt
                )
                texto_respuesta = respuesta.text.strip()

                # Limpiar posible markdown de bloque de código
                texto_respuesta = re.sub(r'^```(?:json)?\s*', '', texto_respuesta)
                texto_respuesta = re.sub(r'\s*```$', '', texto_respuesta)

                resultado = json.loads(texto_respuesta)
                return _validar_coherencia_analisis(resultado)

            except json.JSONDecodeError as e:
                return {
                    "error": f"Error al interpretar la respuesta de Gemini: {str(e)}",
                    "resumen": texto_respuesta[:500] if texto_respuesta else "No se pudo obtener respuesta",
                    "hallazgos": [],
                    "recomendaciones": [],
                    "graficas": []
                }
            except Exception as e:
                error_str = str(e)
                es_503 = "503" in error_str or "UNAVAILABLE" in error_str.upper()
                es_conexion = any(s in error_str.lower() for s in ("disconnected", "connection", "timeout", "timed out"))
                es_429 = "429" in error_str or "RESOURCE_EXHAUSTED" in error_str.upper()
                es_404 = "404" in error_str or "NOT_FOUND" in error_str.upper()
                es_no_encontrado = "no se encuentra" in error_str.lower()

                if es_429:
                    ultimo_error = error_str
                    # Si es el último intento de este modelo, pasar al siguiente
                    if intento >= max_intentos_por_modelo:
                        break
                    # Si no, esperar y reintentar
                    match = re.search(r'retry in ([\d.]+)s', error_str, re.IGNORECASE)
                    espera = float(match.group(1)) + 2 if match else 10
                    time.sleep(espera)
                    continue

                if es_404 or es_no_encontrado:
                    # Modelo no encontrado/deprecado, saltar al siguiente
                    ultimo_error = f"Modelo {modelo} no disponible. Probando siguiente..."
                    break

                if (es_503 or es_conexion) and intento < max_intentos_por_modelo:
                    espera = 2 ** intento
                    time.sleep(espera)
                    continue

                if es_503 or es_conexion:
                    ultimo_error = "Gemini está experimentando alta demanda en este momento."
                    break

                # Error desconocido: devolver inmediatamente
                return {
                    "error": f"Error al conectar con Gemini: {error_str}",
                    "hallazgos": [],
                    "recomendaciones": [],
                    "graficas": []
                }

    # Todos los modelos gratuitos fallaron por cuota.
    # Intentar con API Key de pago si está configurada (respaldo final)
    api_key_pago = os.getenv("GEMINI_API_KEY_PAGO", "").strip()
    if api_key_pago and api_key_pago != api_key:
        try:
            client_pago = genai.Client(api_key=api_key_pago)
            respuesta = client_pago.models.generate_content(
                model=model_principal,
                contents=prompt
            )
            texto_respuesta = respuesta.text.strip()
            texto_respuesta = re.sub(r'^```(?:json)?\s*', '', texto_respuesta)
            texto_respuesta = re.sub(r'\s*```$', '', texto_respuesta)
            resultado = json.loads(texto_respuesta)
            return _validar_coherencia_analisis(resultado)
        except Exception:
            pass

    return {
        "error": (
            "Todos los modelos gratuitos de Gemini han alcanzado su limite de cuota. "
            "Se reintentara automaticamente en unos segundos."
        ),
        "hallazgos": [],
        "recomendaciones": [],
        "graficas": []
    }


# ── Ruta para la página de análisis con IA ──
@informe_bp.route("/analizar_informe/<int:id_informe>")
def analizar_informe(id_informe):
    """Página de análisis del documento con Gemini."""
    if "usuario" not in session:
        return redirect(url_for("home"))
    perfil = session.get("perfil_activo")
    if perfil not in ("fiscalizador", "admin"):
        return acceso_no_autorizado()

    conn = None
    cursor = None

    try:
        conn = conexion()
        cursor = conn.cursor(dictionary=True)

        if perfil == "admin":
            cursor.execute("""
                SELECT i.*, g.fecha_guardia
                FROM informes i
                INNER JOIN guardias g ON i.id_guardia = g.id_guardia
                WHERE i.id_informe = %s
                AND i.estado = 'activo'
            """, (id_informe,))
        else:
            cursor.execute("""
                SELECT i.*, g.fecha_guardia
                FROM informes i
                INNER JOIN guardias g ON i.id_guardia = g.id_guardia
                WHERE i.id_informe = %s
                AND i.id_usuario = %s
                AND i.estado = 'activo'
            """, (id_informe, session["id_usuario"]))
        informe = cursor.fetchone()

        if not informe:
            return no_encontrado("Informe no encontrado")

        return render_template("analizar_informe.html", informe=informe)

    finally:
        if cursor is not None: cursor.close()
        if conn is not None: conn.close()


@informe_bp.route("/analizar_informe/<int:id_informe>/api")
def analizar_informe_api(id_informe):
    """API que extrae texto, llama a Gemini y devuelve JSON.
    Usa cache en BD: solo llama a Gemini si no hay resultado previo o si se fuerza (?force=1)."""

    # Verificar sesión y perfil activo
    if "usuario" not in session:
        return jsonify({"error": "No autorizado"}), 401
    perfil = session.get("perfil_activo")
    if perfil not in ("fiscalizador", "admin"):
        return jsonify({"error": "Acceso no autorizado"}), 403

    forzar = request.args.get("force") == "1"

    conn = None
    cursor = None

    try:
        conn = conexion()
        cursor = conn.cursor(dictionary=True)

        if perfil == "admin":
            cursor.execute("""
                SELECT i.*, g.fecha_guardia
                FROM informes i
                INNER JOIN guardias g ON i.id_guardia = g.id_guardia
                WHERE i.id_informe = %s
                AND i.estado = 'activo'
            """, (id_informe,))
        else:
            cursor.execute("""
                SELECT i.*, g.fecha_guardia
                FROM informes i
                INNER JOIN guardias g ON i.id_guardia = g.id_guardia
                WHERE i.id_informe = %s
                AND i.id_usuario = %s
            AND i.estado = 'activo'
        """, (id_informe, session["id_usuario"]))
        informe = cursor.fetchone()

        if not informe:
            return jsonify({"error": "Informe no encontrado"}), 404

        # Servir desde caché si existe y no se fuerza re-análisis
        if not forzar and informe.get("resultado_analisis"):
            import json as _json
            cache = informe["resultado_analisis"]
            if isinstance(cache, str):
                cache = _json.loads(cache)
            return jsonify(cache)

        ruta = informe["ruta_archivo"]
        if not os.path.exists(ruta):
            return jsonify({"error": "El archivo no existe en el servidor"}), 404

        # Extraer texto del archivo
        texto = _extraer_texto_archivo(informe)

        if not texto:
            return jsonify({"error": "No se pudo extraer texto del archivo"}), 400

        if texto.startswith("[Error"):
            return jsonify({"error": texto}), 500

        # Analizar con Gemini
        resultado = _analizar_con_gemini(
            texto,
            informe["titulo"] or "",
            informe["descripcion"] or ""
        )

        # Guardar en caché si el análisis fue exitoso (sin error)
        if not resultado.get("error"):
            try:
                cursor.execute(
                    "UPDATE informes SET resultado_analisis = %s WHERE id_informe = %s",
                    (json.dumps(resultado, ensure_ascii=False), id_informe)
                )
                conn.commit()
            except Exception as e:
                print(f"[ADVERTENCIA] No se pudo guardar el analisis en cache (informe {id_informe}): {e}")

        return jsonify(resultado)

    finally:
        if cursor is not None: cursor.close()
        if conn is not None: conn.close()


@informe_bp.route("/analizar_informe/<int:id_informe>/ppt", methods=["POST"])
def generar_ppt_analisis(id_informe):
    """Genera una presentación PPTX profesional con los resultados del análisis."""

    # Validar sesión y datos de entrada
    if "usuario" not in session:
        return jsonify({"error": "No autorizado"}), 401
    perfil = session.get("perfil_activo")
    if perfil not in ("fiscalizador", "admin"):
        return jsonify({"error": "Acceso no autorizado"}), 403

    data = request.get_json(silent=True)
    if not data or not isinstance(data, dict):
        return jsonify({"error": "Datos no proporcionados o formato invalido"}), 400

    conn = None
    cursor = None

    try:
        conn = conexion()
        cursor = conn.cursor(dictionary=True)

        if perfil == "admin":
            cursor.execute("""
                SELECT i.titulo, i.nombre_archivo, i.tipo_archivo,
                       g.fecha_guardia,
                       CONCAT(u.nombre, ' ', u.apellidos) AS fiscalizador
                FROM informes i
                INNER JOIN guardias g ON i.id_guardia = g.id_guardia
                INNER JOIN usuarios u ON i.id_usuario = u.id_usuario
                WHERE i.id_informe = %s AND i.estado = 'activo'
            """, (id_informe,))
        else:
            cursor.execute("""
                SELECT i.titulo, i.nombre_archivo, i.tipo_archivo,
                       g.fecha_guardia,
                       CONCAT(u.nombre, ' ', u.apellidos) AS fiscalizador
                FROM informes i
                INNER JOIN guardias g ON i.id_guardia = g.id_guardia
                INNER JOIN usuarios u ON i.id_usuario = u.id_usuario
                WHERE i.id_informe = %s AND i.id_usuario = %s AND i.estado = 'activo'
            """, (id_informe, session["id_usuario"]))
        informe = cursor.fetchone()

        if not informe:
            return jsonify({"error": "Informe no encontrado"}), 404

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Paleta de colores corporativos para la presentación
        C_BG_DARK   = RGBColor(0x08, 0x0E, 0x1A)
        C_PRIMARY   = RGBColor(0x1E, 0x3A, 0x5F)
        C_ACCENT    = RGBColor(0x63, 0x66, 0xF1)
        C_ACCENT2   = RGBColor(0x8B, 0x5C, 0xF6)
        C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
        C_LIGHT_BG  = RGBColor(0xF8, 0xFA, 0xFC)
        C_CARD_BG   = RGBColor(0xFF, 0xFF, 0xFF)
        C_TEXT      = RGBColor(0x1E, 0x29, 0x3B)
        C_SUBTLE    = RGBColor(0x94, 0xA3, 0xB8)
        C_GOLD      = RGBColor(0xD9, 0x77, 0x06)
        C_GREEN     = RGBColor(0x05, 0x96, 0x69)
        C_RED       = RGBColor(0xDC, 0x26, 0x26)
        C_BLUE      = RGBColor(0x3B, 0x82, 0xF6)
        C_ORANGE    = RGBColor(0xEA, 0x58, 0x0C)
        C_TEAL      = RGBColor(0x0D, 0x94, 0x8B)
        C_BORDER    = RGBColor(0xE2, 0xE8, 0xF0)

        paleta_modulos = [C_BLUE, C_GREEN, C_ACCENT2, C_ORANGE, C_TEAL, C_ACCENT, C_GOLD]

        from datetime import date as date_type

        # ── Funciones helper para construir la presentación ──

        def _set_bg(slide, color):
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = color

        def _rounded_rect(slide, x, y, w, h, fill_color=C_CARD_BG, border_color=None, radius=0.15):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            if border_color:
                shape.line.color.rgb = border_color
                shape.line.width = Pt(0.8)
            else:
                shape.line.fill.background()
            return shape

        def _text_box(slide, x, y, w, h, text, size=13, color=C_TEXT, bold=False, align=PP_ALIGN.LEFT, word_wrap=True):
            tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf = tb.text_frame; tf.word_wrap = word_wrap
            p = tf.paragraphs[0]; p.alignment = align
            r = p.add_run(); r.text = str(text)
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
            return tb, tf

        def _multi_text(slide, x, y, w, lines, size=13, color=C_TEXT, bold_first=False, line_spacing=Pt(18)):
            tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(len(lines) * 0.4 + 0.3))
            tf = tb.text_frame; tf.word_wrap = True
            for i, txt in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                r = p.add_run(); r.text = str(txt)
                r.font.size = Pt(size); r.font.color.rgb = color
                r.font.bold = (bold_first and i == 0)
                p.line_spacing = line_spacing
            return tb

        def _progress_bar(slide, x, y, w, pct, color, label=""):
            h = 0.32
            track = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
            )
            track.fill.solid(); track.fill.fore_color.rgb = RGBColor(0xE8, 0xEC, 0xF1)
            track.line.fill.background()

            pct_clamped = max(0, min(100, float(pct) if pct else 0)) / 100.0
            fill_w = max(0.1, (w - 0.08) * pct_clamped) if pct_clamped > 0 else 0
            if fill_w > 0:
                fill_bar = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.04), Inches(y + 0.04),
                    Inches(fill_w), Inches(h - 0.08)
                )
                fill_bar.fill.solid(); fill_bar.fill.fore_color.rgb = color
                fill_bar.line.fill.background()

            if label:
                _text_box(slide, x + 0.1, y + 0.02, w - 0.2, 0.28, label, 9, C_WHITE, True, PP_ALIGN.LEFT)

            pct_label = f"{float(pct):.0f}%" if pct else "0%"
            _text_box(slide, x + w - 0.7, y + 0.02, 0.65, 0.28, pct_label, 9, C_TEXT, True, PP_ALIGN.RIGHT)

        def _add_header(slide, title, accent_color, subtitle=None, module_num=None):
            bar_h = 1.25
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(bar_h)
            )
            bar.fill.solid(); bar.fill.fore_color.rgb = C_BG_DARK; bar.line.fill.background()

            accent_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(bar_h), Inches(13.333), Inches(0.05)
            )
            accent_line.fill.solid(); accent_line.fill.fore_color.rgb = accent_color; accent_line.line.fill.background()

            if module_num:
                badge = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, Inches(0.6), Inches(0.22), Inches(0.8), Inches(0.8)
                )
                badge.fill.solid(); badge.fill.fore_color.rgb = accent_color; badge.line.fill.background()
                btf = badge.text_frame; btf.paragraphs[0].alignment = PP_ALIGN.CENTER
                br = btf.paragraphs[0].add_run(); br.text = str(module_num)
                br.font.size = Pt(18); br.font.color.rgb = C_WHITE; br.font.bold = True

            title_x = 1.7 if module_num else 0.8
            tb = slide.shapes.add_textbox(Inches(title_x), Inches(0.15), Inches(10.5), Inches(1.0))
            tf = tb.text_frame; tf.word_wrap = True
            tp = tf.paragraphs[0]; tr = tp.add_run()
            tr.text = title; tr.font.size = Pt(24); tr.font.color.rgb = C_WHITE; tr.font.bold = True
            if subtitle:
                sp = tf.add_paragraph(); sr = sp.add_run()
                sr.text = subtitle; sr.font.size = Pt(12); sr.font.color.rgb = C_SUBTLE

        def _add_kpi(slide, x, y, w, value, label, color, icon_text="■"):
            card = _rounded_rect(slide, x, y, w, 1.35, C_CARD_BG, C_BORDER)
            top = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x + 0.15), Inches(y + 0.08), Inches(w - 0.3), Inches(0.05)
            )
            top.fill.solid(); top.fill.fore_color.rgb = color; top.line.fill.background()

            _text_box(slide, x + 0.2, y + 0.25, w - 0.4, 0.55, str(value), 26, color, True, PP_ALIGN.CENTER)
            _text_box(slide, x + 0.2, y + 0.85, w - 0.4, 0.4, label, 10, C_SUBTLE, False, PP_ALIGN.CENTER)

        def _workshop_card(slide, x, y, w, taller, ejecutadas, completadas, eficacia, color):
            card = _rounded_rect(slide, x, y, w, 1.45, C_CARD_BG, C_BORDER)

            left_stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.08), Inches(0.06), Inches(1.29)
            )
            left_stripe.fill.solid(); left_stripe.fill.fore_color.rgb = color; left_stripe.line.fill.background()

            _text_box(slide, x + 0.25, y + 0.12, w - 0.5, 0.35, taller, 14, C_TEXT, True)
            _text_box(slide, x + 0.25, y + 0.5, w - 0.5, 0.3, f"{ejecutadas} ejecutadas | {completadas} completadas", 10, C_SUBTLE)

            _progress_bar(slide, x + 0.25, y + 0.92, w - 0.5, eficacia, color)

        def _section_title(slide, x, y, text, color=C_TEXT, size=15):
            _text_box(slide, x, y, 5, 0.35, text, size, color, True)

        # ═══════════════════════════════
        # S1 — PORTADA (título, metadatos y total de actividades)
        # ═══════════════════════════════
        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(s1, C_BG_DARK)

        s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08)
        ).fill.solid(); s1.shapes[-1].fill.fore_color.rgb = C_ACCENT; s1.shapes[-1].line.fill.background()

        for i, c in enumerate([C_ACCENT, C_ACCENT2, C_BLUE, C_GREEN]):
            s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.8 + i*0.45), Inches(2.0), Inches(0.1), Inches(0.1)
            ).fill.solid(); s1.shapes[-1].fill.fore_color.rgb = c; s1.shapes[-1].line.fill.background()

        logo_shape = s1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(2.4), Inches(0.75), Inches(0.75)
        )
        logo_shape.fill.solid(); logo_shape.fill.fore_color.rgb = C_ACCENT; logo_shape.line.fill.background()
        lt = logo_shape.text_frame; lt.paragraphs[0].alignment = PP_ALIGN.CENTER
        lr = lt.paragraphs[0].add_run(); lr.text = "AI"; lr.font.size = Pt(22)
        lr.font.color.rgb = C_WHITE; lr.font.bold = True

        _text_box(s1, 2.75, 2.45, 8.5, 0.4, "REPORTE DE ANÁLISIS DE GUARDIA", 15, C_ACCENT, True)
        _text_box(s1, 2.75, 2.85, 8.5, 0.6, informe["titulo"] or "Informe sin título", 30, C_WHITE, True)

        s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.75), Inches(3.6), Inches(2.5), Inches(0.035)
        ).fill.solid(); s1.shapes[-1].fill.fore_color.rgb = C_ACCENT; s1.shapes[-1].line.fill.background()

        meta_lines = [
            f"Fecha de Guardia: {informe.get('fecha_guardia', 'N/A')}",
            f"Fiscalizador: {informe.get('fiscalizador', 'N/A')}",
            f"Documento: {informe.get('nombre_archivo', '')} ({informe.get('tipo_archivo', '').upper()})"
        ]
        _multi_text(s1, 2.75, 3.9, 8.5, meta_lines, 13, C_SUBTLE)

        big_num = data.get("modulo1_resumen_ejecutivo", {}).get("total_actividades", "—")
        _text_box(s1, 9.5, 1.8, 3.2, 0.7, str(big_num), 60, C_ACCENT, True, PP_ALIGN.CENTER)
        _text_box(s1, 9.5, 2.4, 3.2, 0.3, "ACTIVIDADES\nEJECUTADAS", 10, C_SUBTLE, False, PP_ALIGN.CENTER)

        bottom = s1.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.5), Inches(13.333), Inches(1.0)
        )
        bottom.fill.solid(); bottom.fill.fore_color.rgb = C_PRIMARY; bottom.line.fill.background()
        bt = bottom.text_frame; bt.word_wrap = True
        bp = bt.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = f"SIGGO — Integrated Management System for Guards and Operations  |  {date_type.today().strftime('%d/%m/%Y')}"
        br.font.size = Pt(11); br.font.color.rgb = C_SUBTLE
        bp2 = bt.add_paragraph(); bp2.alignment = PP_ALIGN.CENTER
        br2 = bp2.add_run(); br2.text = "Análisis potenciado por Inteligencia Artificial"
        br2.font.size = Pt(10); br2.font.color.rgb = C_SUBTLE

        # ═══════════════════════════════
        # S2 — MÓDULO 1: RESUMEN EJECUTIVO (KPIs principales y resumen textual)
        # ═══════════════════════════════
        mod1 = data.get("modulo1_resumen_ejecutivo") or {}
        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(s2, C_LIGHT_BG)
        _add_header(s2, "Resumen Ejecutivo de Guardia", C_ACCENT, "MÓDULO 1 — Indicadores consolidados", 1)

        kpis_config = [
            (mod1.get("total_actividades", "—"), "Total Actividades\nEjecutadas", C_BLUE),
            (mod1.get("total_produccion_recuperada_bopd", "—"), "Prod. Recuperada\nCNP (BOPD)", C_GREEN),
            (mod1.get("total_trabajos_criticos", "—"), "Trabajos\nCríticos", C_GOLD),
            (mod1.get("eficacia_general", "—"), "Eficacia\nGeneral", C_ACCENT2),
        ]
        for i, (val, lbl, clr) in enumerate(kpis_config):
            _add_kpi(s2, 0.6 + i*3.15, 1.7, 2.95, val, lbl, clr)

        resumen = data.get("resumen") or "No se generó resumen."
        res_box = _rounded_rect(s2, 0.6, 3.4, 12.1, 3.8, C_CARD_BG, C_BORDER)
        lines_resumen = resumen[:1000].split("\n")[:14]
        _multi_text(s2, 0.9, 3.6, 11.5, lines_resumen, 12, C_TEXT, False, Pt(19))

        # ═══════════════════════════════
        # S3 — MÓDULO 2: EFICACIA POR TALLER (tarjetas con barra de progreso)
        # ═══════════════════════════════
        eficacia = data.get("modulo2_eficacia_taller") or []
        if eficacia:
            s3 = prs.slides.add_slide(prs.slide_layouts[6])
            _set_bg(s3, C_LIGHT_BG)
            _add_header(s3, "Control de Actividades y Eficacia por Taller", C_BLUE, "MÓDULO 2 — Desempeño operativo por especialidad", 2)

            cols = min(3, len(eficacia))
            for i, row in enumerate(eficacia):
                col = i % cols
                fila = i // cols
                cx = 0.6 + col * 4.1
                cy = 1.7 + fila * 1.7
                taller_name = row.get("taller", f"Taller {i+1}")
                color = paleta_modulos[i % len(paleta_modulos)]
                ej = row.get("ejecutadas", 0)
                comp = row.get("completadas", 0)
                ef = row.get("eficacia", 0)
                _workshop_card(s3, cx, cy, 3.85, taller_name, ej, comp, ef, color)

        # ═══════════════════════════════
        # S4 — MÓDULO 3: MATRIZ DE TRABAJOS IMPORTANTES (críticos / amarillos)
        # ═══════════════════════════════
        importantes = data.get("modulo3_trabajos_importantes") or []
        if importantes:
            s4 = prs.slides.add_slide(prs.slide_layouts[6])
            _set_bg(s4, C_LIGHT_BG)
            _add_header(s4, "Matriz de Trabajos Importantes", C_GOLD, "MÓDULO 3 — Trabajos críticos / celdas amarillas", 3)

            alert = _rounded_rect(s4, 0.6, 1.6, 12.1, 0.5, RGBColor(0xFF, 0xFB, 0xEB), C_GOLD)
            _text_box(s4, 0.9, 1.68, 11.5, 0.35,
                f"⚠  {len(importantes)} trabajos críticos identificados que requieren atención prioritaria",
                13, C_GOLD, True, PP_ALIGN.CENTER)

            rows_per_page = min(7, len(importantes))
            for i in range(0, len(importantes), rows_per_page):
                if i > 0:
                    s4 = prs.slides.add_slide(prs.slide_layouts[6])
                    _set_bg(s4, C_LIGHT_BG)
                    _add_header(s4, "Matriz de Trabajos Importantes (cont.)", C_GOLD, None, 3)

                chunk = importantes[i:i+rows_per_page]
                for j, row in enumerate(chunk):
                    cy = 2.4 + j * 0.72
                    taller = row.get("taller", ""); pozo = row.get("pozo", "")
                    falla = row.get("falla", ""); solucion = row.get("solucion", "")

                    num_circle = s4.shapes.add_shape(
                        MSO_SHAPE.OVAL, Inches(0.7), Inches(cy + 0.05), Inches(0.42), Inches(0.42)
                    )
                    num_circle.fill.solid(); num_circle.fill.fore_color.rgb = C_GOLD
                    num_circle.line.fill.background()
                    ntf = num_circle.text_frame; ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    nr = ntf.paragraphs[0].add_run(); nr.text = str(i + j + 1)
                    nr.font.size = Pt(12); nr.font.color.rgb = C_WHITE; nr.font.bold = True

                    card = _rounded_rect(s4, 1.3, cy, 11.4, 0.62, C_CARD_BG, C_BORDER)
                    left_dot = s4.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, Inches(1.3), Inches(cy + 0.04), Inches(0.05), Inches(0.54)
                    )
                    left_dot.fill.solid(); left_dot.fill.fore_color.rgb = C_GOLD; left_dot.line.fill.background()

                    _text_box(s4, 1.55, cy + 0.02, 2.0, 0.28, taller, 9, C_GOLD, True)
                    _text_box(s4, 1.55, cy + 0.28, 2.0, 0.28, f"Pozo: {pozo}", 9, C_SUBTLE)
                    _text_box(s4, 3.6, cy + 0.04, 4.2, 0.52, f"Falla: {falla[:120]}", 10, C_TEXT)
                    _text_box(s4, 8.0, cy + 0.04, 4.5, 0.52, f"Solución: {solucion[:120]}", 10, C_TEXT)

        # ═══════════════════════════════
        # S5 — MÓDULO 4: RESUMEN OPERATIVO (producción recuperada y clasificación)
        # ═══════════════════════════════
        ops = data.get("modulo4_resumen_operativo") or {}
        bopd_data = ops.get("bopd_por_taller") or {}
        tipos_data = ops.get("trabajos_por_tipo") or {}
        if bopd_data or tipos_data:
            s5 = prs.slides.add_slide(prs.slide_layouts[6])
            _set_bg(s5, C_LIGHT_BG)
            _add_header(s5, "Resumen Operativo Adicional", C_GREEN, "MÓDULO 4 — Producción recuperada y clasificación de trabajos", 4)

            if bopd_data:
                total_bopd = sum(float(v) for v in bopd_data.values() if v)
                _add_kpi(s5, 0.6, 1.7, 5.8, f"{total_bopd:.1f}", "Total CNP Recuperado (BOPD)", C_GREEN)
                cy = 3.4
                max_bopd = max(float(v) for v in bopd_data.values()) if bopd_data else 1
                for k, v in bopd_data.items():
                    pct_b = (float(v) / max_bopd * 100) if max_bopd > 0 else 0
                    _text_box(s5, 0.8, cy, 2.2, 0.28, str(k), 11, C_TEXT, True)
                    _text_box(s5, 2.6, cy, 1.0, 0.28, f"{v} bopd", 11, C_TEXT)
                    _progress_bar(s5, 3.8, cy + 0.02, 2.8, pct_b, C_GREEN, "")
                    cy += 0.45

            start_x = 7.2 if bopd_data else 0.6
            tip_ty = 1.7
            if tipos_data:
                _section_title(s5, start_x, tip_ty, "TRABAJOS POR TIPO")
                tip_ty += 0.45
                type_colors = {"CNP": C_GREEN, "MC": C_GOLD, "PV": C_BLUE, "SOP": C_ACCENT, "OP": C_ACCENT2, "SUS": C_RED}
                max_t = max(int(v) for v in tipos_data.values()) if tipos_data else 1
                for tipo, cant in tipos_data.items():
                    pct_t = (int(cant) / max_t * 100) if max_t > 0 else 0
                    tclr = type_colors.get(str(tipo).upper(), C_ACCENT)
                    _text_box(s5, start_x + 0.2, tip_ty, 1.5, 0.28, str(tipo), 13, tclr, True)
                    _text_box(s5, start_x + 1.8, tip_ty, 1.0, 0.28, str(cant), 13, C_TEXT, True)
                    _progress_bar(s5, start_x + 2.9, tip_ty + 0.02, 3.3, pct_t, tclr, "")
                    tip_ty += 0.45

        # ═══════════════════════════════
        # S6 — VISUALIZACIÓN DE KPIs POR GRÁFICAS (barras horizontales)
        # ═══════════════════════════════
        graficas = data.get("graficas") or []
        if graficas:
            s6 = prs.slides.add_slide(prs.slide_layouts[6])
            _set_bg(s6, C_LIGHT_BG)
            _add_header(s6, "Visualización de Datos", C_ACCENT2, "KPIs del análisis operativo", 5)

            for gidx, g in enumerate(graficas):
                titulo = g.get("titulo", f"Gráfica {gidx+1}")
                labels = g.get("labels", [])
                datasets = g.get("datasets", [])
                if not datasets or not labels:
                    continue

                cy = 1.7 + gidx * 1.8
                if cy > 6.0:
                    s6 = prs.slides.add_slide(prs.slide_layouts[6])
                    _set_bg(s6, C_LIGHT_BG)
                    cy = 0.6

                _section_title(s6, 0.8, cy, titulo, C_PRIMARY, 14)
                cy += 0.42

                ds = datasets[0]
                valores = ds.get("data", [])
                max_val = max(float(v) for v in valores) if valores else 1
                paleta = ds.get("backgroundColor", []) or [
                    "#3b82f6", "#22c55e", "#f59e0b", "#8b5cf6", "#ef4444", "#06b6d4", "#f97316", "#64748b"
                ]

                for li, (label, val) in enumerate(zip(labels, valores)):
                    pct = (float(val) / max_val * 100) if max_val > 0 else 0
                    bar_color = RGBColor(
                        int(paleta[li % len(paleta)][1:3], 16),
                        int(paleta[li % len(paleta)][3:5], 16),
                        int(paleta[li % len(paleta)][5:7], 16)
                    )
                    _text_box(s6, 0.8, cy, 2.4, 0.26, str(label)[:20], 10, C_TEXT, True)
                    _text_box(s6, 3.2, cy, 0.8, 0.26, str(val), 10, C_TEXT)
                    _progress_bar(s6, 4.2, cy + 0.01, 8.3, pct, bar_color, "")
                    cy += 0.35

        # ═══════════════════════════════
        # S7 — MÓDULO 5: PENDIENTES (tareas en espera con punto de estado)
        # ═══════════════════════════════
        pendientes = data.get("modulo5_pendientes") or []
        if pendientes:
            s7 = prs.slides.add_slide(prs.slide_layouts[6])
            _set_bg(s7, C_LIGHT_BG)
            _add_header(s7, "Reporte de Actividades Pendientes", C_RED, "MÓDULO 5 — Tareas en espera registradas", 5)

            alert2 = _rounded_rect(s7, 0.6, 1.6, 12.1, 0.5, RGBColor(0xFE, 0xF2, 0xF2), C_RED)
            _text_box(s7, 0.9, 1.68, 11.5, 0.35,
                f"⚠  {len(pendientes)} actividades pendientes que requieren seguimiento inmediato",
                13, C_RED, True, PP_ALIGN.CENTER)

            for i, row in enumerate(pendientes):
                if i >= 10: break
                cy = 2.4 + i * 0.5
                taller = row.get("taller", ""); pozo = row.get("pozo", "")
                req = row.get("requerimiento", ""); estado = row.get("estado", "Pendiente")

                status_color = C_GOLD if str(estado).lower() in ("pendiente", "en proceso") else C_RED
                status_dot = s7.shapes.add_shape(
                    MSO_SHAPE.OVAL, Inches(0.75), Inches(cy + 0.06), Inches(0.22), Inches(0.22)
                )
                status_dot.fill.solid(); status_dot.fill.fore_color.rgb = status_color
                status_dot.line.fill.background()

                _text_box(s7, 1.15, cy, 2.0, 0.22, f"{taller} | {pozo}", 10, C_TEXT, True)
                _text_box(s7, 3.4, cy, 6.5, 0.22, str(req)[:100], 10, C_SUBTLE)
                _text_box(s7, 10.4, cy, 2.3, 0.22, str(estado).upper(), 10, status_color, True)

        # ═══════════════════════════════
        # S9 — DIAPOSITIVA DE CIERRE
        # ═══════════════════════════════
        s_end = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(s_end, C_BG_DARK)

        end_line = s_end.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(2.8), Inches(2.3), Inches(0.04)
        )
        end_line.fill.solid(); end_line.fill.fore_color.rgb = C_ACCENT; end_line.line.fill.background()

        _text_box(s_end, 1.5, 3.1, 10.3, 1.0, "FIN DEL REPORTE", 44, C_WHITE, True, PP_ALIGN.CENTER)
        _text_box(s_end, 1.5, 4.1, 10.3, 0.5,
            "Análisis generado automáticamente con Inteligencia Artificial",
            14, C_SUBTLE, False, PP_ALIGN.CENTER)

        end_footer = s_end.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(10.333), Inches(1.0))
        eft = end_footer.text_frame; eft.word_wrap = True
        ep5 = eft.paragraphs[0]; ep5.alignment = PP_ALIGN.CENTER
        er5 = ep5.add_run()
        er5.text = "SIGGO — Integrated Management System for Guards and Operations"
        er5.font.size = Pt(13); er5.font.color.rgb = C_SUBTLE; er5.font.italic = True

        # ── Guardar presentación en memoria y enviar como descarga ──
        output = BytesIO()
        prs.save(output)
        output.seek(0)

        nombre_archivo = secure_filename(informe.get("titulo", "analisis") or "analisis")
        nombre_ppt = f"Analisis_IA_{nombre_archivo}.pptx"

        return send_file(
            output,
            as_attachment=True,
            download_name=nombre_ppt,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except Exception as e:
        return jsonify({"error": f"Error al generar PPT: {str(e)}"}), 500

    finally:
        if cursor is not None: cursor.close()
        if conn is not None: conn.close()
